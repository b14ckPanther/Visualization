import PyPDF2
import os
import json

def extract_pdf(path):
    try:
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n---PAGE---\n"
            return text
    except Exception as e:
        return f"ERROR: {e}"

def extract_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except:
        return "DOCX extraction failed"

def extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n---SLIDE---\n"
        return text
    except:
        return "PPTX extraction failed"

results = {}
dirs = ["Lectures", "Practice", "Syllabus"]

for d in dirs:
    dirpath = os.path.join("/Users/zangeel/Desktop/VisualizationExam", d)
    for f in sorted(os.listdir(dirpath)):
        fpath = os.path.join(dirpath, f)
        if f.endswith('.pdf'):
            results[f"{d}/{f}"] = extract_pdf(fpath)
        elif f.endswith('.docx'):
            results[f"{d}/{f}"] = extract_docx(fpath)
        elif f.endswith('.pptx'):
            results[f"{d}/{f}"] = extract_pptx(fpath)

# Write each to a separate text file for easier reading
outdir = "/Users/zangeel/Desktop/VisualizationExam/extracted"
os.makedirs(outdir, exist_ok=True)

for key, text in results.items():
    safe_name = key.replace("/", "__").replace(" ", "_") + ".txt"
    with open(os.path.join(outdir, safe_name), 'w', encoding='utf-8') as f:
        f.write(text[:50000])  # cap at 50k chars per file
    print(f"Extracted: {key} -> {len(text)} chars")
